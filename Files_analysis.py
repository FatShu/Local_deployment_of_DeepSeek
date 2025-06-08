import os
import json
import requests
import pandas as pd
from typing import Optional, List

# 全局配置
API_BASE_URL = "http://localhost:11434"

# 尝试导入docx库，用于处理Word文档
try:
    from docx import Document
except ImportError:
    Document = None

# 尝试导入textract库，用于处理.doc文件
try:
    import textract
except ImportError:
    textract = None

# 尝试导入python-pptx库，用于处理PowerPoint文档
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def get_database_dir() -> str:
    """获取Database文件夹的路径"""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(script_dir, "Database")

def get_all_analyzable_files() -> List[str]:
    """获取Database文件夹下所有可分析文件的名称"""
    database_dir = get_database_dir()
    if not os.path.exists(database_dir):
        print(f"错误: Database文件夹 '{database_dir}' 不存在")
        return []
    
    supported_extensions = ['.txt', '.xlsx', '.docx', '.doc', '.pptx', '.ppt']
    return [f for f in os.listdir(database_dir) 
            if os.path.isfile(os.path.join(database_dir, f)) 
            and any(f.lower().endswith(ext) for ext in supported_extensions)]

def read_file_content(file_path: str) -> Optional[str]:
    """读取文件内容，支持txt、xlsx、docx、doc、pptx、ppt格式"""
    try:
        if file_path.lower().endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
                
        elif file_path.lower().endswith('.xlsx'):
            # 读取Excel文件
            df = pd.read_excel(file_path)
            # 将数据框转换为CSV格式的字符串，保留表头
            return df.to_csv(sep='\t', na_rep='nan')
            
        elif file_path.lower().endswith('.docx'):
            if Document is None:
                print("错误: 要处理docx文件，需要安装python-docx库")
                print("输入pip install python-docx进行安装")
                return None
                
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return '\n'.join(full_text)
            
        elif file_path.lower().endswith('.doc'):
            if textract is None:
                print("错误: 要处理doc文件，需要安装textract库")
                print("输入pip install textract进行安装")
                return None
                
            try:
                text = textract.process(file_path).decode('utf-8')
                return text
            except Exception as e:
                print(f"错误: 无法读取doc文件 '{file_path}': {e}")
                return None
                
        elif file_path.lower().endswith(('.pptx', '.ppt')):
            if Presentation is None:
                print("错误: 要处理PowerPoint文件，需要安装python-pptx库")
                print("输入pip install python-pptx进行安装")
                return None
                
            if file_path.lower().endswith('.pptx'):
                prs = Presentation(file_path)
                full_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            full_text.append(shape.text)
                return '\n'.join(full_text)
            else:  # .ppt
                print("警告: 对.ppt文件的支持有限，建议转换为.pptx格式")
                if textract is None:
                    print("错误: 要处理.ppt文件，需要安装textract库")
                    print("输入pip install textract进行安装")
                    return None
                try:
                    text = textract.process(file_path).decode('utf-8')
                    return text
                except Exception as e:
                    print(f"错误: 无法读取ppt文件 '{file_path}': {e}")
                    return None
                
        else:
            print(f"错误: 不支持的文件格式: {file_path}")
            return None
            
    except FileNotFoundError:
        print(f"错误: 文件 '{file_path}' 未找到")
        return None
    except Exception as e:
        print(f"错误: 无法读取文件 '{file_path}': {e}")
        return None

def list_available_models() -> List[str]:
    """获取Ollama服务上可用的模型列表"""
    global API_BASE_URL
    try:
        print(f"正在连接Ollama: {API_BASE_URL}...")
        url = f"{API_BASE_URL}/api/tags"
        response = requests.get(url, timeout=5)
        response.raise_for_status()
        
        result = response.json()
        models = result.get("models", [])
        if not models:
            print("警告: Ollama中未找到任何模型")
        return [model["name"] for model in models]
    except requests.exceptions.ConnectionError:
        print("错误: 无法连接到Ollama。请确认:")
        print("  1. Ollama服务正在运行")
        print(f"  2. Ollama API地址正确 (当前监听地址: {API_BASE_URL})")
        return []
    except requests.exceptions.Timeout:
        print("错误: 连接Ollama服务超时。请检查服务是否正常运行。")
        return []
    except requests.exceptions.HTTPError as e:
        print(f"错误: Ollama服务返回HTTP错误 {e.response.status_code}")
        return []
    except json.JSONDecodeError:
        print("错误: 无法解析Ollama服务的响应")
        return []
    except Exception as e:
        print(f"发生未知错误: {e}")
        return []

def ask_ollama(prompt: str, model: str = "deepseek-r1:1.5b", system: str = "你是一个智能助手，擅长分析文本内容。") -> Optional[str]:
    """向Ollama API发送请求并获取流式响应"""
    global API_BASE_URL
    try:
        url = f"{API_BASE_URL}/api/generate"
        headers = {"Content-Type": "application/json"}
        data = {
            "model": model,
            "prompt": prompt,
            "system": system,
            "stream": True  # 启用流式响应
        }
        
        print(f"\n {model} 正在思考")
        print("=" * 50)
        
        full_response = ""
        with requests.post(url, headers=headers, json=data, stream=True, timeout=300) as response:
            response.raise_for_status()
            
            # 处理流式响应
            for line in response.iter_lines():
                if line:
                    try:
                        # 解析JSON行
                        chunk = json.loads(line.decode('utf-8'))
                        
                        # 检查是否有内容
                        if 'response' in chunk:
                            text_chunk = chunk['response']
                            full_response += text_chunk
                            
                            # 实时输出
                            print(text_chunk, end='', flush=True)
                            
                        # 检查是否结束
                        if chunk.get('done', False):
                            break
                            
                    except json.JSONDecodeError:
                        print("\n错误: 无法解析")
                        return None
        
        print("\n" + "=" * 50)
        return full_response
        
    except requests.exceptions.ConnectionError:
        print("\n错误: 无法连接到Ollama服务。请检查服务是否正在运行。")
        return None
    except requests.exceptions.Timeout:
        print("\n错误: 请求超时。请检查模型是否正在加载或资源是否充足。")
        return None
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 404:
            print(f"\n错误: 模型 '{model}' 不存在。请确保:")
            print("  1. 该模型已安装到Ollama服务中")
            print("  2. 可以使用 'ollama pull [模型名称]' 安装模型")
        else:
            print(f"\n错误: Ollama服务返回HTTP错误 {e.response.status_code}")
        return None
    except Exception as e:
        print(f"\n发生未知错误: {e}")
        return None

def analyze_file_with_ollama(file_name: str, analysis_type: str, model: str = "deepseek-r1:1.5b") -> None:
    """读取并使用Ollama分析单个文件"""
    database_dir = get_database_dir()
    file_path = os.path.join(database_dir, file_name)
    
    print(f"\n正在分析文件: {file_name}")
    content = read_file_content(file_path)
    if not content:
        return
    
    # 根据文件类型调整系统提示
    if file_name.lower().endswith('.xlsx'):
        system_prompt = "你是一个数据分析专家，擅长分析和解读表格数据。"
    elif file_name.lower().endswith(('.doc', '.docx')):
        system_prompt = "你是一个文档分析专家，擅长理解和总结文档内容。"
    elif file_name.lower().endswith(('.pptx', '.ppt')):
        system_prompt = "你是一个演示文稿分析专家，擅长总结幻灯片内容和提取关键点。"
    else:
        system_prompt = "你是一个智能助手，擅长分析文本内容。"
    
    # 根据文件大小调整提示词
    file_size = os.path.getsize(file_path)
    if file_size > 1024 * 1024:  # 1MB以上
        print("警告: 文件较大，可能会超出模型上下文限制")
        prompt = f"以下是文件内容的一部分：\n\n{content[:20000]}\n\n请{analysis_type}这段内容，重点关注关键信息。"
    else:
        prompt = f"请{analysis_type}以下文件内容：\n\n{content}"
    
    # 获取并返回完整响应
    response = ask_ollama(prompt, model, system_prompt)
    return response

def analyze_all_files_interactive(model: str = "deepseek-r1:1.5b") -> None:
    """交互式分析Database目录下的所有支持的文件"""
    files = get_all_analyzable_files()
    
    if not files:
        print("未找到可分析的文件(txt、xlsx、docx、doc、pptx、ppt)")
        return
    
    print(f"找到 {len(files)} 个可分析文件")
    
    for i, file_name in enumerate(files, 1):
        print(f"\n文件 {i}/{len(files)}: {file_name}")
        
        while True:
            analysis_type = input("\n请输入要对该文件执行的操作 (例如: 总结, 提取要点, 分析数据, 查找特定信息, 退出): ").strip()
            
            if not analysis_type:
                print("操作不能为空，请重新输入")
                continue
            
            if analysis_type.lower() == "退出":
                print(f"跳过文件: {file_name}")
                break
                
            analyze_file_with_ollama(file_name, analysis_type, model)
            
            # 询问是否继续处理此文件
            continue_processing = input("\n是否要对该文件执行其他操作? (y/n): ").strip().lower()
            if continue_processing != 'y':
                break

def select_model_from_user() -> str:
    """获取用户选择的模型"""
    global API_BASE_URL
    models = list_available_models()
    
    if not models:
        print("\n错误: 未找到可用模型。请确认:")
        print("  1. Ollama服务正在运行")
        print("  2. 至少有一个模型已安装到Ollama服务中")
        print("  3. 可以使用 'ollama pull [模型名称]' 安装模型")
        
        # 询问是否更改了端口
        change_port = input("\n是否更改了Ollama的监听端口? (y/n): ").strip().lower()
        if change_port != 'y':
            print("程序将退出。请启动Ollama服务并安装模型后再试。")
            exit(1)
            
        # 获取新端口并尝试连接
        new_port = input("请输入Ollama服务的端口号 (默认: 11434): ").strip()
        if not new_port:
            new_port = "11434"
            
        try:
            port_num = int(new_port)
            if not (1 <= port_num <= 65535):
                print("错误: 端口号必须在1-65535之间")
                exit(1)
        except ValueError:
            print("错误: 端口号必须是数字")
            exit(1)
            
        API_BASE_URL = f"http://localhost:{new_port}"
        print(f"已更新Ollama API地址为: {API_BASE_URL}")
        
        # 重新尝试获取模型列表
        print("\n正在使用新端口重新连接Ollama服务...")
        models = list_available_models()
        
        if not models:
            print("\n错误: 即使更改端口后，仍无法连接到Ollama服务或未找到任何模型。")
            print("请确认:")
            print("  1. Ollama服务正在运行")
            print(f"  2. Ollama服务正在监听端口 {new_port}")
            print("  3. 至少有一个模型已安装到Ollama服务中")
            exit(1)
    
    print("\n可用模型:")
    for i, model in enumerate(models, 1):
        print(f"{i}. {model}")
    
    while True:
        choice = input("\n请输入要使用的模型编号或名称 (默认: 1): ").strip()
        
        if not choice:
            return models[0]
        
        try:
            index = int(choice) - 1
            if 0 <= index < len(models):
                return models[index]
            else:
                print("编号超出范围，请重新输入")
        except ValueError:
            if choice in models:
                return choice
            else:
                print("无效的模型名称，请重新输入")

if __name__ == "__main__":
    # 确保安装了必要的库
    required_libraries = {
        'pandas': 'pandas',
        'openpyxl': 'openpyxl'
    }
    
    # 检查docx支持
    if Document is None:
        required_libraries['docx'] = 'python-docx'
    
    # 检查doc支持
    if textract is None:
        required_libraries['textract'] = 'textract'
    
    # 检查PowerPoint支持
    if Presentation is None:
        required_libraries['pptx'] = 'python-pptx'
    
    missing_libraries = []
    for lib_name, pkg_name in required_libraries.items():
        try:
            __import__(lib_name)
        except ImportError:
            missing_libraries.append(pkg_name)
    
    if missing_libraries:
        print(f"错误: 需要安装以下库以支持所有文件类型: {', '.join(missing_libraries)}")
        print(f"输入pip install {' '.join(missing_libraries)}进行安装")
        exit(1)
    
    # 让用户选择模型
    selected_model = select_model_from_user()
    print(f"已选择模型: {selected_model}")
    
    # 交互式分析文件
    analyze_all_files_interactive(selected_model)    
